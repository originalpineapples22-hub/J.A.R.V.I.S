# -*- coding: utf-8 -*-
"""Office suite: real PowerPoint, Word and Excel files, styled, downloadable."""
from pathlib import Path
from . import tool
from .. import memory
from ..config import FILES_DIR

DARK = (5, 11, 26)
ACC = (63, 169, 255)
LIGHT = (234, 244, 255)


def _safe(name, ext):
    n = Path(name or f"document{ext}").name
    return n if n.lower().endswith(ext) else n + ext


@tool("make_presentation", "Create a styled PowerPoint (.pptx). Provide a title and a list of slides, each with a title, bullet points and optional speaker notes.",
      {"name": "file name", "title": "deck title", "subtitle": "optional", "slides": "[{\"title\": str, \"bullets\": [str], \"notes\": str}]", "theme": "dark|light"},
      agent="Coding Agent")
def make_presentation(args, ctx):
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    dark = (args.get("theme", "dark") != "light")
    bg, fg, acc = (DARK, LIGHT, ACC) if dark else ((255, 255, 255), (20, 30, 50), (20, 110, 200))
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)

    def paint(slide):
        f = slide.background.fill
        f.solid()
        f.fore_color.rgb = RGBColor(*bg)

    def text(slide, x, y, w, h, s, size, bold=False, color=fg):
        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = s
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = RGBColor(*color)
        return tf

    s = prs.slides.add_slide(prs.slide_layouts[6])
    paint(s)
    text(s, 0.8, 2.4, 11.5, 1.5, args.get("title", "Untitled"), 44, True, acc)
    if args.get("subtitle"):
        text(s, 0.8, 3.9, 11.5, 1, args["subtitle"], 22)
    for sl in args.get("slides", []) or []:
        s = prs.slides.add_slide(prs.slide_layouts[6])
        paint(s)
        text(s, 0.7, 0.5, 12, 1, sl.get("title", ""), 32, True, acc)
        line = s.shapes.add_shape(1, Inches(0.7), Inches(1.35), Inches(2.5), Inches(0.05))
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(*acc)
        line.line.fill.background()
        tf = text(s, 0.9, 1.7, 11.5, 5, "", 20)
        bullets = sl.get("bullets", []) or []
        for i, b in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = "•  " + str(b)
            p.font.size = Pt(20)
            p.font.color.rgb = RGBColor(*fg)
            p.space_after = Pt(10)
        if sl.get("notes"):
            s.notes_slide.notes_text_frame.text = str(sl["notes"])
    name = _safe(args.get("name") or args.get("title", "presentation"), ".pptx")
    prs.save(str(FILES_DIR / name))
    memory.add_event("file", f"Presentation created: {name}")
    return f"Presentation saved as {name} ({1 + len(args.get('slides', []) or [])} slides). Download it from Files."


@tool("make_document", "Create a Word document (.docx) with a title and sections of headings and paragraphs.",
      {"name": "file name", "title": "document title", "sections": "[{\"heading\": str, \"paragraphs\": [str], \"bullets\": [str]}]"},
      agent="Coding Agent")
def make_document(args, ctx):
    from docx import Document
    from docx.shared import Pt, RGBColor
    d = Document()
    d.add_heading(args.get("title", "Document"), 0)
    for sec in args.get("sections", []) or []:
        if sec.get("heading"):
            d.add_heading(sec["heading"], 1)
        for p in sec.get("paragraphs", []) or []:
            d.add_paragraph(str(p))
        for b in sec.get("bullets", []) or []:
            d.add_paragraph(str(b), style="List Bullet")
    name = _safe(args.get("name") or args.get("title", "document"), ".docx")
    d.save(str(FILES_DIR / name))
    memory.add_event("file", f"Document created: {name}")
    return f"Document saved as {name}. Download it from Files."


@tool("make_spreadsheet", "Create an Excel workbook (.xlsx). Provide sheets with headers and rows; numbers are kept numeric and formulas (strings starting with '=') work.",
      {"name": "file name", "sheets": "[{\"name\": str, \"headers\": [str], \"rows\": [[...]]}]"}, agent="Coding Agent")
def make_spreadsheet(args, ctx):
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill
    wb = Workbook()
    wb.remove(wb.active)
    for sh in args.get("sheets", []) or [{"name": "Sheet1", "headers": [], "rows": []}]:
        ws = wb.create_sheet(title=str(sh.get("name", "Sheet"))[:31])
        headers = sh.get("headers", []) or []
        if headers:
            ws.append(headers)
            for c in ws[1]:
                c.font = Font(bold=True, color="FFFFFF")
                c.fill = PatternFill("solid", fgColor="1F4E9C")
        for row in sh.get("rows", []) or []:
            ws.append(list(row))
        for col in ws.columns:
            ws.column_dimensions[col[0].column_letter].width = min(40, max(10, max(len(str(c.value or "")) for c in col) + 2))
    name = _safe(args.get("name", "workbook"), ".xlsx")
    wb.save(str(FILES_DIR / name))
    memory.add_event("file", f"Spreadsheet created: {name}")
    return f"Spreadsheet saved as {name}. Download it from Files."
